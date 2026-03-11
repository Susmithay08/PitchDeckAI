"""
PPTX + PDF export for PitchDeck AI.
Supports multiple color themes.
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── THEMES ──────────────────────────────────────────────────────────────────
THEMES = {
    "midnight": {
        "name": "Midnight Executive",
        "bg_dark":    "1E2761",
        "bg_light":   "F8F9FC",
        "accent":     "4FC3F7",
        "accent2":    "FFD54F",
        "text_light": "FFFFFF",
        "text_dark":  "1E2761",
        "card_bg":    "243380",
        "muted":      "90CAF9",
    },
    "coral": {
        "name": "Coral Energy",
        "bg_dark":    "2F3C7E",
        "bg_light":   "FFF8F5",
        "accent":     "F96167",
        "accent2":    "F9E795",
        "text_light": "FFFFFF",
        "text_dark":  "2F3C7E",
        "card_bg":    "3D4F9F",
        "muted":      "F9A8AE",
    },
    "forest": {
        "name": "Forest & Moss",
        "bg_dark":    "1B4332",
        "bg_light":   "F5F9F5",
        "accent":     "52B788",
        "accent2":    "B7E4C7",
        "text_light": "FFFFFF",
        "text_dark":  "1B4332",
        "card_bg":    "2D6A4F",
        "muted":      "95D5B2",
    },
    "berry": {
        "name": "Berry & Cream",
        "bg_dark":    "6D2E46",
        "bg_light":   "FFF8F7",
        "accent":     "E07A5F",
        "accent2":    "F2CC8F",
        "text_light": "FFFFFF",
        "text_dark":  "6D2E46",
        "card_bg":    "8B3A5A",
        "muted":      "EDAFB8",
    },
    "slate": {
        "name": "Charcoal Minimal",
        "bg_dark":    "1C1C2E",
        "bg_light":   "F7F7F8",
        "accent":     "7C3AED",
        "accent2":    "A78BFA",
        "text_light": "FFFFFF",
        "text_dark":  "1C1C2E",
        "card_bg":    "2D2D44",
        "muted":      "9CA3AF",
    },
    "ocean": {
        "name": "Ocean Deep",
        "bg_dark":    "065A82",
        "bg_light":   "F0F7FF",
        "accent":     "1DD3B0",
        "accent2":    "AFFC41",
        "text_light": "FFFFFF",
        "text_dark":  "065A82",
        "card_bg":    "1C7293",
        "muted":      "72EFDD",
    },
}


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def _add_bg(slide, color_hex):
    """Fill slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def _txt(frame, text, size, bold=False, color="FFFFFF", align=PP_ALIGN.LEFT, italic=False):
    tf = frame.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)


def _add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    return shape


def _add_textbox(slide, text, x, y, w, h, size, bold=False, color="FFFFFF",
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)
    return txBox


def _add_textbox_multiline(slide, lines, x, y, w, h, size, color, bold=False, line_gap_pt=6):
    """Add multiple lines as separate paragraphs in one textbox."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_gap_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(color)


# ─── SLIDE BUILDERS ──────────────────────────────────────────────────────────

def slide_cover(prs, t, data, pitch):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full dark bg
    _add_bg(slide, t["bg_dark"])
    # Accent bar top
    _add_rect(slide, 0, 0, 10, 0.06, t["accent"])
    # Big startup name
    _add_textbox(slide, pitch.get("startup_name",""), 0.7, 1.2, 8.6, 1.6, 52, bold=True,
                 color=t["text_light"], align=PP_ALIGN.LEFT)
    # Tagline italic
    _add_textbox(slide, pitch.get("tagline",""), 0.7, 2.9, 8.6, 0.7, 22, italic=True,
                 color=t["accent"], align=PP_ALIGN.LEFT)
    # Divider
    _add_rect(slide, 0.7, 3.7, 1.5, 0.05, t["accent2"])
    # Body
    _add_textbox(slide, data.get("body",""), 0.7, 3.9, 8.2, 1.0, 14,
                 color=t["muted"], align=PP_ALIGN.LEFT)
    # Slide label bottom right
    _add_textbox(slide, "01 / COVER", 8.0, 5.1, 1.7, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_problem(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_light"])
    # Left accent bar
    _add_rect(slide, 0, 0, 0.08, 5.625, t["accent"])
    # Title
    _add_textbox(slide, data.get("title","The Problem"), 0.3, 0.35, 5, 0.7, 30, bold=True,
                 color=t["text_dark"])
    # Big stat card
    _add_rect(slide, 6.5, 0.3, 3.2, 1.6, t["bg_dark"])
    _add_textbox(slide, data.get("stat",""), 6.6, 0.45, 3.0, 0.8, 26, bold=True,
                 color=t["accent"], align=PP_ALIGN.CENTER)
    _add_textbox(slide, "MARKET PAIN", 6.6, 1.15, 3.0, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.CENTER)
    # Body
    _add_textbox(slide, data.get("body",""), 0.3, 1.2, 5.8, 1.0, 13,
                 color=t["text_dark"])
    # Bullets
    bullets = data.get("bullets", [])
    for i, b in enumerate(bullets[:3]):
        yp = 2.4 + i * 0.85
        _add_rect(slide, 0.3, yp, 0.06, 0.55, t["accent"])
        _add_textbox(slide, f"0{i+1}", 0.5, yp, 0.4, 0.5, 11, bold=True, color=t["accent"])
        _add_textbox(slide, b, 0.95, yp, 8.7, 0.6, 12, color=t["text_dark"])
    _add_textbox(slide, "02 / PROBLEM", 8.0, 5.1, 1.7, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_solution(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_dark"])
    _add_rect(slide, 0, 0, 10, 0.06, t["accent2"])
    _add_textbox(slide, data.get("title","Our Solution"), 0.7, 0.4, 8, 0.7, 30, bold=True,
                 color=t["text_light"])
    _add_textbox(slide, data.get("body",""), 0.7, 1.2, 8.6, 1.0, 13,
                 color=t["muted"])
    bullets = data.get("bullets", [])
    for i, b in enumerate(bullets[:3]):
        yp = 2.35 + i * 0.9
        _add_rect(slide, 0.7, yp, 0.5, 0.55, t["accent"])
        _add_textbox(slide, "✓", 0.7, yp, 0.5, 0.55, 14, bold=True,
                     color=t["bg_dark"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, b, 1.35, yp + 0.05, 8.3, 0.6, 13, color=t["text_light"])
    _add_textbox(slide, "03 / SOLUTION", 7.8, 5.1, 2.0, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_market(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_light"])
    _add_rect(slide, 0, 0, 0.08, 5.625, t["accent2"])
    _add_textbox(slide, data.get("title","Market Opportunity"), 0.3, 0.35, 7, 0.65, 30, bold=True,
                 color=t["text_dark"])
    _add_textbox(slide, data.get("body",""), 0.3, 1.1, 9.4, 0.8, 12, color=t["text_dark"])
    # TAM/SAM/SOM cards
    cards = [("TAM", data.get("tam",""), "Total Addressable Market"),
             ("SAM", data.get("sam",""), "Serviceable Market"),
             ("SOM", data.get("som",""), "Obtainable Market")]
    colors = [t["bg_dark"], t["accent"], t["card_bg"]]
    for i, (label, val, desc) in enumerate(cards):
        xp = 0.3 + i * 3.2
        _add_rect(slide, xp, 2.1, 3.0, 2.8, colors[i])
        _add_textbox(slide, label, xp + 0.15, 2.2, 2.7, 0.4, 10, bold=True,
                     color=t["muted"] if i == 0 else t["bg_dark"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, val, xp + 0.1, 2.65, 2.8, 0.9, 30, bold=True,
                     color=t["accent2"] if i == 0 else t["bg_dark"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, desc, xp + 0.1, 3.65, 2.8, 0.8, 10,
                     color=t["muted"] if i == 0 else t["card_bg"], align=PP_ALIGN.CENTER)
    _add_textbox(slide, "04 / MARKET", 8.0, 5.1, 1.8, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_product(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_dark"])
    _add_rect(slide, 0, 0, 10, 0.06, t["accent"])
    _add_textbox(slide, data.get("title","How It Works"), 0.7, 0.3, 8, 0.7, 30, bold=True,
                 color=t["text_light"])
    steps = data.get("steps", [])
    step_colors = [t["accent"], t["accent2"], t["card_bg"]]
    for i, step in enumerate(steps[:3]):
        xp = 0.5 + i * 3.1
        _add_rect(slide, xp, 1.2, 2.85, 3.5, step_colors[i % len(step_colors)])
        num_color = t["bg_dark"] if i < 2 else t["accent"]
        _add_textbox(slide, step.get("step",""), xp + 0.15, 1.3, 2.5, 0.8, 36, bold=True,
                     color=num_color, align=PP_ALIGN.LEFT)
        _add_textbox(slide, step.get("title",""), xp + 0.15, 2.15, 2.55, 0.65, 14, bold=True,
                     color=t["bg_dark"] if i < 2 else t["text_light"])
        _add_textbox(slide, step.get("desc",""), xp + 0.15, 2.85, 2.55, 1.65, 11,
                     color=t["bg_dark"] if i < 2 else t["muted"])
    _add_textbox(slide, "05 / PRODUCT", 7.8, 5.1, 2.0, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_traction(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_light"])
    _add_rect(slide, 0, 0, 0.08, 5.625, t["accent"])
    _add_textbox(slide, data.get("title","Traction"), 0.3, 0.35, 8, 0.65, 30, bold=True,
                 color=t["text_dark"])
    metrics = data.get("metrics", [])
    bg_alts = [t["bg_dark"], t["card_bg"], t["accent"], t["bg_dark"]]
    for i, m in enumerate(metrics[:4]):
        xp = 0.3 + (i % 2) * 4.8
        yp = 1.15 + (i // 2) * 1.65
        _add_rect(slide, xp, yp, 4.5, 1.45, bg_alts[i])
        _add_textbox(slide, m.get("value",""), xp + 0.2, yp + 0.05, 4.1, 0.8, 28, bold=True,
                     color=t["accent"] if i % 2 == 0 else t["accent2"])
        _add_textbox(slide, m.get("label",""), xp + 0.2, yp + 0.85, 4.1, 0.45, 11,
                     color=t["muted"])
    _add_textbox(slide, data.get("body",""), 0.3, 4.6, 9.4, 0.6, 11, color=t["text_dark"])
    _add_textbox(slide, "06 / TRACTION", 7.8, 5.1, 2.0, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_business_model(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_dark"])
    _add_rect(slide, 0, 0, 10, 0.06, t["accent2"])
    _add_textbox(slide, data.get("title","Business Model"), 0.7, 0.3, 8, 0.65, 30, bold=True,
                 color=t["text_light"])
    streams = data.get("revenue_streams", [])
    colors = [t["accent"], t["accent2"], t["card_bg"]]
    for i, rs in enumerate(streams[:3]):
        yp = 1.2 + i * 1.2
        _add_rect(slide, 0.7, yp, 0.5, 0.9, colors[i % len(colors)])
        _add_textbox(slide, rs.get("name",""), 1.35, yp, 4.0, 0.5, 14, bold=True,
                     color=t["text_light"])
        _add_textbox(slide, rs.get("desc",""), 1.35, yp + 0.42, 5.5, 0.55, 11,
                     color=t["muted"])
        _add_textbox(slide, rs.get("pct",""), 8.8, yp + 0.05, 1.0, 0.65, 22, bold=True,
                     color=t["accent"], align=PP_ALIGN.RIGHT)
    _add_textbox(slide, "07 / BUSINESS MODEL", 7.1, 5.1, 2.7, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_competition(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_light"])
    _add_rect(slide, 0, 0, 0.08, 5.625, t["accent2"])
    _add_textbox(slide, data.get("title","Competition"), 0.3, 0.35, 6, 0.65, 30, bold=True,
                 color=t["text_dark"])
    _add_textbox(slide, data.get("body",""), 0.3, 1.1, 5.8, 1.2, 13, color=t["text_dark"])
    _add_textbox(slide, "OUR ADVANTAGES", 6.4, 0.35, 3.3, 0.4, 10, bold=True,
                 color=t["accent"])
    advs = data.get("advantages", [])
    for i, a in enumerate(advs[:4]):
        yp = 0.9 + i * 1.0
        _add_rect(slide, 6.4, yp, 0.35, 0.35, t["accent"])
        _add_textbox(slide, "★", 6.4, yp, 0.35, 0.35, 12, color=t["bg_dark"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, a, 6.9, yp, 2.8, 0.7, 12, color=t["text_dark"])
    _add_textbox(slide, "08 / COMPETITION", 7.5, 5.1, 2.3, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_team(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_dark"])
    _add_rect(slide, 0, 0, 10, 0.06, t["accent"])
    _add_textbox(slide, data.get("title","The Team"), 0.7, 0.3, 8, 0.65, 30, bold=True,
                 color=t["text_light"])
    members = data.get("members", [])
    member_colors = [t["accent"], t["accent2"], t["card_bg"]]
    for i, m in enumerate(members[:3]):
        xp = 0.5 + i * 3.1
        # Avatar circle (square)
        _add_rect(slide, xp, 1.2, 0.75, 0.75, member_colors[i % len(member_colors)])
        name = m.get("name","")
        _add_textbox(slide, name[0] if name else "?", xp, 1.2, 0.75, 0.75, 18, bold=True,
                     color=t["bg_dark"], align=PP_ALIGN.CENTER)
        _add_textbox(slide, name, xp, 2.1, 2.8, 0.55, 15, bold=True, color=t["text_light"])
        _add_textbox(slide, m.get("role",""), xp, 2.6, 2.8, 0.4, 11,
                     color=member_colors[i % len(member_colors)])
        _add_textbox(slide, m.get("background",""), xp, 3.05, 2.8, 1.4, 11,
                     color=t["muted"])
    _add_textbox(slide, "09 / TEAM", 8.2, 5.1, 1.6, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


def slide_ask(prs, t, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, t["bg_light"])
    _add_rect(slide, 0, 0, 0.08, 5.625, t["accent"])
    _add_textbox(slide, data.get("title","The Ask"), 0.3, 0.35, 5, 0.65, 30, bold=True,
                 color=t["text_dark"])
    # Amount
    _add_textbox(slide, data.get("amount",""), 0.3, 1.1, 5, 1.2, 54, bold=True,
                 color=t["accent"])
    _add_textbox(slide, data.get("round","") + " ROUND", 0.3, 2.35, 5, 0.4, 12,
                 color=t["text_dark"])
    # Use of funds bars
    funds = data.get("use_of_funds", [])
    bar_colors = [t["bg_dark"], t["accent"], t["card_bg"], t["accent2"]]
    _add_textbox(slide, "USE OF FUNDS", 5.7, 0.8, 4.0, 0.4, 10, bold=True, color=t["accent"])
    for i, f in enumerate(funds[:4]):
        yp = 1.35 + i * 0.9
        pct = f.get("pct", 0)
        bar_w = (pct / 100) * 3.8
        _add_rect(slide, 5.7, yp, 3.8, 0.3, "E5E7EB")
        _add_rect(slide, 5.7, yp, max(bar_w, 0.05), 0.3, bar_colors[i % len(bar_colors)])
        _add_textbox(slide, f.get("category",""), 5.7, yp + 0.32, 2.5, 0.3, 11,
                     color=t["text_dark"])
        _add_textbox(slide, f"{pct}%", 9.0, yp + 0.32, 0.7, 0.3, 11, bold=True,
                     color=t["accent"], align=PP_ALIGN.RIGHT)
    # Vision quote
    if data.get("vision"):
        _add_rect(slide, 0.3, 4.1, 9.4, 1.15, t["bg_dark"])
        _add_textbox(slide, f'"{data.get("vision","")}"', 0.5, 4.15, 9.0, 1.0, 13,
                     italic=True, color=t["text_light"])
    _add_textbox(slide, "10 / THE ASK", 8.2, 5.1, 1.6, 0.4, 9,
                 color=t["muted"], align=PP_ALIGN.RIGHT)


SLIDE_BUILDERS = {
    "cover": slide_cover,
    "problem": slide_problem,
    "solution": slide_solution,
    "market": slide_market,
    "product": slide_product,
    "traction": slide_traction,
    "business_model": slide_business_model,
    "competition": slide_competition,
    "team": slide_team,
    "ask": slide_ask,
}


def build_pptx(pitch_data: dict, theme_key: str = "midnight") -> bytes:
    """Build a PPTX from pitch data. Returns bytes."""
    t = THEMES.get(theme_key, THEMES["midnight"])
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_data in pitch_data.get("slides", []):
        sid = slide_data.get("id", "")
        builder = SLIDE_BUILDERS.get(sid)
        if not builder:
            continue
        if sid == "cover":
            builder(prs, t, slide_data, pitch_data)
        else:
            builder(prs, t, slide_data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def build_pdf(pitch_data: dict, theme_key: str = "midnight") -> bytes:
    """Build a fully styled, themed PDF pitch deck using reportlab directly."""
    return _build_pdf_reportlab(pitch_data, theme_key)


def _build_pdf_reportlab(pitch_data: dict, theme_key: str) -> bytes:
    """
    Fully themed PDF using reportlab. Matches the theme colors exactly.
    16:9 landscape format (792 x 446 pt).
    """
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor
    from reportlab.lib.utils import simpleSplit

    t = THEMES.get(theme_key, THEMES["midnight"])
    W, H = 792, 446   # 16:9 landscape

    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=(W, H))
    c.setTitle(pitch_data.get("startup_name", "PitchDeck AI"))

    def col(hx):
        hx = hx.lstrip("#")
        return HexColor(f"#{hx}")

    def rect(x, y, w, h, fill_hex, stroke=False):
        c.setFillColor(col(fill_hex))
        c.rect(x, y, w, h, fill=1, stroke=1 if stroke else 0)

    def text(s, x, y, size, color_hex, bold=False, italic=False, align="left", max_w=None):
        if not s: return
        font = "Helvetica-BoldOblique" if bold and italic else "Helvetica-Bold" if bold else "Helvetica-Oblique" if italic else "Helvetica"
        c.setFont(font, size)
        c.setFillColor(col(color_hex))
        if align == "right":
            c.drawRightString(x, y, str(s))
        elif align == "center":
            c.drawCentredString(x, y, str(s))
        else:
            if max_w:
                lines = simpleSplit(str(s), font, size, max_w)
                for line in lines[:3]:
                    c.drawString(x, y, line)
                    y -= size + 2
            else:
                c.drawString(x, y, str(s))

    def wrapped(s, x, y, max_w, size, color_hex, bold=False, max_lines=4):
        if not s: return y
        font = "Helvetica-Bold" if bold else "Helvetica"
        c.setFont(font, size)
        c.setFillColor(col(color_hex))
        lines = simpleSplit(str(s), font, size, max_w)
        for line in lines[:max_lines]:
            c.drawString(x, y, line)
            y -= size + 3
        return y

    def slide_bg(dark=True):
        rect(0, 0, W, H, t["bg_dark"] if dark else t["bg_light"])
        # Accent top bar
        rect(0, H - 5, W, 5, t["accent"])
        # Slide number bottom right
        c.setFont("Helvetica", 8)
        c.setFillColor(col(t["muted"]))

    def footer(num, label):
        c.setFont("Helvetica", 8)
        c.setFillColor(col(t["muted"]))
        c.drawRightString(W - 20, 14, f"{str(num).zfill(2)} / {label.upper()}")

    slides = pitch_data.get("slides", [])
    name = pitch_data.get("startup_name", "")
    tagline = pitch_data.get("tagline", "")
    total = len(slides)

    for idx, s in enumerate(slides):
        sid = s.get("id", "")
        use_dark = sid in ("cover", "solution", "product", "business_model", "team")

        # ── Background ──
        slide_bg(use_dark)
        footer(idx + 1, s.get("title", ""))

        tl = t["text_light"]
        td = t["text_dark"]
        tc = tl if use_dark else td
        mt = t["muted"]
        ac = t["accent"]
        ac2 = t["accent2"]

        if sid == "cover":
            # Big name
            text(name, 44, H - 80, 38, tl, bold=True)
            # Accent rule
            rect(44, H - 96, 80, 3, ac)
            # Tagline
            text(tagline, 44, H - 118, 16, ac, italic=True)
            # Body
            wrapped(s.get("body",""), 44, H - 150, W * 0.55, 11, mt, max_lines=4)
            # Big decorative letter top-right
            c.setFont("Helvetica-Bold", 220)
            c.setFillColor(HexColor(f"#{t['card_bg']}"))
            c.drawRightString(W - 20, H - 220, name[0].upper() if name else "P")

        elif sid == "problem":
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            # Stat card top right
            rect(W * 0.6, H - 130, W * 0.37, 80, t["card_bg"])
            text(s.get("stat",""), W * 0.615, H - 80, 22, ac, bold=True)
            text("MARKET PAIN", W * 0.615, H - 100, 8, mt)
            # Body text
            wrapped(s.get("body",""), 44, H - 85, W * 0.52, 11, tc, max_lines=3)
            # Bullets
            y = H - 150
            for j, b in enumerate(s.get("bullets",[])[:3]):
                rect(44, y - 2, 3, 18, ac)
                text(f"0{j+1}", 52, y, 9, ac, bold=True)
                wrapped(b, 70, y, W * 0.5, 10, tc, max_lines=2)
                y -= 32

        elif sid == "solution":
            text(s.get("title",""), 44, H - 56, 24, tl, bold=True)
            wrapped(s.get("body",""), 44, H - 88, W * 0.88, 11, mt, max_lines=3)
            y = H - 150
            for b in s.get("bullets",[])[:3]:
                rect(44, y - 2, 24, 20, ac)
                text("✓", 50, y, 11, t["bg_dark"], bold=True)
                wrapped(b, 76, y, W * 0.82, 11, tl, max_lines=2)
                y -= 36

        elif sid == "market":
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            wrapped(s.get("body",""), 44, H - 84, W * 0.88, 11, tc, max_lines=2)
            cards = [("TAM", s.get("tam",""), t["bg_dark"]), ("SAM", s.get("sam",""), t["card_bg"]), ("SOM", s.get("som",""), ac)]
            cx = 44
            for lbl, val, bg in cards:
                rect(cx, H - 310, 218, 130, bg)
                text(lbl, cx + 12, H - 210, 8, mt)
                text(val, cx + 12, H - 246, 26, ac2, bold=True)
                cx += 232

        elif sid == "product":
            text(s.get("title",""), 44, H - 56, 24, tl, bold=True)
            step_bgs = [t["accent"], t["accent2"], t["card_bg"]]
            sx = 44
            for j, step in enumerate(s.get("steps",[])[:3]):
                bg = step_bgs[j % 3]
                rect(sx, H - 340, 218, 200, bg)
                txt_c = t["bg_dark"] if j < 2 else tl
                text(step.get("step",""), sx + 14, H - 112, 32, txt_c, bold=True)
                wrapped(step.get("title",""), sx + 14, H - 158, 190, 12, txt_c, bold=True, max_lines=2)
                wrapped(step.get("desc",""), sx + 14, H - 190, 190, 10, txt_c if j > 1 else t["bg_dark"], max_lines=4)
                sx += 232

        elif sid == "traction":
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            metrics = s.get("metrics",[])[:4]
            positions = [(44, H - 230), (W//2 + 10, H - 230), (44, H - 360), (W//2 + 10, H - 360)]
            mcolors = [t["bg_dark"], t["card_bg"], t["card_bg"], t["bg_dark"]]
            for j, m in enumerate(metrics):
                mx, my = positions[j]
                rect(mx, my, W//2 - 54, 110, mcolors[j])
                text(m.get("value",""), mx + 14, my + 70, 26, ac if j % 2 == 0 else ac2, bold=True)
                text(m.get("label",""), mx + 14, my + 46, 9, mt)
            wrapped(s.get("body",""), 44, 35, W - 88, 10, tc, max_lines=2)

        elif sid == "business_model":
            text(s.get("title",""), 44, H - 56, 24, tl, bold=True)
            y = H - 110
            stream_cols = [ac, ac2, t["card_bg"]]
            for j, rs in enumerate(s.get("revenue_streams",[])[:3]):
                rect(44, y - 20, 6, 36, stream_cols[j % 3])
                text(rs.get("name",""), 58, y, 13, tl, bold=True)
                wrapped(rs.get("desc",""), 58, y - 18, W * 0.65, 10, mt, max_lines=1)
                text(rs.get("pct",""), W - 44, y, 18, ac, bold=True, align="right")
                y -= 70

        elif sid == "competition":
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            wrapped(s.get("body",""), 44, H - 90, W * 0.45, 11, tc, max_lines=4)
            text("OUR ADVANTAGES", W * 0.53, H - 80, 9, ac, bold=True)
            y = H - 110
            for a in s.get("advantages",[])[:4]:
                rect(W * 0.53, y - 2, 18, 16, ac)
                text("★", W * 0.53 + 3, y, 9, td)
                wrapped(a, W * 0.53 + 24, y, W * 0.42, 10, tc, max_lines=2)
                y -= 36

        elif sid == "team":
            text(s.get("title",""), 44, H - 56, 24, tl, bold=True)
            mem_cols = [ac, ac2, t["card_bg"]]
            mx = 44
            for j, m in enumerate(s.get("members",[])[:3]):
                mc = mem_cols[j % 3]
                rect(mx, H - 120, 46, 46, mc)
                init = m.get("name","?")[0]
                text(init, mx + 13, H - 82, 22, td, bold=True)
                text(m.get("name",""), mx, H - 138, 13, tl, bold=True)
                text(m.get("role",""), mx, H - 154, 9, mc)
                wrapped(m.get("background",""), mx, H - 178, 210, 9, mt, max_lines=4)
                mx += 232

        elif sid == "ask":
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            text(s.get("amount",""), 44, H - 120, 46, ac, bold=True)
            text((s.get("round","") + " ROUND").upper(), 44, H - 145, 10, tc)
            # Use of funds bars
            text("USE OF FUNDS", W * 0.48, H - 80, 9, ac, bold=True)
            fy = H - 110
            bar_w = W * 0.45
            for f in s.get("use_of_funds",[])[:4]:
                pct = f.get("pct", 0)
                text(f.get("category",""), W * 0.48, fy, 10, tc, bold=True)
                rect(W * 0.48, fy - 14, bar_w, 10, t["card_bg"] if use_dark else "DDDDDD")
                rect(W * 0.48, fy - 14, bar_w * (pct / 100), 10, ac)
                text(f"{pct}%", W * 0.48 + bar_w + 6, fy - 5, 9, ac, bold=True)
                fy -= 40
            # Vision quote box
            if s.get("vision"):
                rect(44, 26, W * 0.38, 44, t["card_bg"])
                wrapped(f'"{s.get("vision","")}"', 54, 60, W * 0.36, 9, tl, max_lines=3)

        else:
            # Generic fallback
            text(s.get("title",""), 44, H - 56, 24, tc, bold=True)
            y = wrapped(s.get("body",""), 44, H - 90, W - 88, 11, tc, max_lines=3)
            y -= 10
            for b in (s.get("bullets") or s.get("advantages") or [])[:5]:
                rect(44, y, 4, 14, ac)
                wrapped(b, 54, y, W - 100, 10, tc, max_lines=2)
                y -= 28

        c.showPage()

    c.save()
    buf.seek(0)
    return buf.read()
